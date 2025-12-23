# ==================== PRODUCTION ENVIRONMENT SETUP ====================
# (Dependencies are handled by Dockerfile)

import os
import io
import re
import ast
import docx
import math
import json
import openai
import base64
import getpass
import zipfile
import openpyxl
import pdfplumber
import pandas as pd
import pytesseract
import uvicorn
from fastapi import FastAPI, HTTPException, BackgroundTasks
from pydantic import BaseModel, Field
from typing import List, Optional
from google.cloud import storage
from tqdm import tqdm
from PIL import Image
from io import StringIO
from pptx import Presentation
from bs4 import BeautifulSoup
from google.cloud import bigquery
from xml.etree import ElementTree as ET
from datetime import datetime, timedelta
from langdetect import detect, DetectorFactory, LangDetectException

import config

# ==================== AUTHENTICATION ====================
# (Production uses Workload Identity / ADC - handled by BigQuery client automatically)
print("🔧 Production environment detected - assuming Workload Identity for GCP access")

DetectorFactory.seed = 0

# Configure Tesseract - production safe approach
if os.environ.get('TESSERACT_PATH'):
    pytesseract.pytesseract.tesseract_cmd = os.environ['TESSERACT_PATH']
else:
    common_paths = [
        r"/usr/bin/tesseract", # Linux/Docker default
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
    ]
    for path in common_paths:
        if os.path.exists(path):
            pytesseract.pytesseract.tesseract_cmd = path
            break

# NON-TRANSLATABLE PHRASES
NON_TRANSLATABLE_PATTERNS = config.NON_TRANSLATABLE_PATTERNS

# Path Configuration Aliases (Aditya's requirement)
DATA_BASE_DIR = config.DATA_BASE_DIR
LOG_DIR = config.LOG_DIR
OUTPUT_DIR = config.OUTPUT_DIR
BENCHMARK_FILE_ID = config.BENCHMARK_FILE_ID
BENCHMARK_LOCAL_PATH = config.BENCHMARK_LOCAL_PATH
FALLBACK_SLA_PATH = config.FALLBACK_SLA_PATH

os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(LOG_DIR, exist_ok=True)
# ==================== API MODELS ====================

class ScopingRequest(BaseModel):
    job_ids: str = Field(..., description="Comma-separated list of job IDs")
    gcs_input_path: str = Field(..., description="GCS path to the folder containing input files (e.g., users/hash/job/timestamp/)")
    instructions: Optional[str] = ""
    ramped_daily_throughput: Optional[float] = None
    ramp_up_days: Optional[int] = None
    translator_pct: Optional[float] = 0.6
    reviewer_pct: Optional[float] = 0.3
    pm_pct: Optional[float] = 0.1

app = FastAPI(title="Lilt Scoping Agent API")

# ==================== CONFIGURATION & DEFAULTS ====================

def download_file_from_google_drive(file_id, destination):
    """
    Downloads a file from Google Drive, handling large file warnings.
    """
    import requests
    URL = "https://docs.google.com/uc?export=download"
    session = requests.Session()

    # Attempt download
    response = session.get(URL, params={'id': file_id}, stream=True)

    # Check for Google's large file download confirmation
    token = None
    for key, value in response.cookies.items():
        if key.startswith('download_warning'):
            token = value
            break

    if token:
        response = session.get(URL, params={'id': file_id, 'confirm': token}, stream=True)

    # Save to disk
    with open(destination, "wb") as f:
        for chunk in response.iter_content(32768):
            if chunk:
                f.write(chunk)

def is_valid_parquet(file_path):
    """Basic check to ensure file is a valid Parquet (starts with 'PAR1')"""
    if not os.path.exists(file_path) or os.path.getsize(file_path) < 4:
        return False
    try:
        with open(file_path, 'rb') as f:
            return f.read(4) == b'PAR1'
    except:
        return False

# ==================== GCS HELPERS ====================

def download_from_gcs(gcs_path: str, local_dir: str):
    """
    Downloads all files from a GCS prefix to a local directory.
    If gcs_path is a local directory, it uses those files directly (for testing).
    """
    if not gcs_path.startswith("gs://"):
        # Local test bypass
        if os.path.isfile(gcs_path):
            import shutil
            dest = os.path.join(local_dir, os.path.basename(gcs_path))
            shutil.copy2(gcs_path, dest)
            print(f"🏠 Local file detected: {gcs_path}")
            return [dest]
        elif os.path.isdir(gcs_path):
            print(f"🏠 Local directory detected: {gcs_path}")
            local_files = []
            for f in os.listdir(gcs_path):
                full_path = os.path.join(gcs_path, f)
                if os.path.isfile(full_path):
                    # Copy to job local_dir to maintain isolation
                    import shutil
                    dest = os.path.join(local_dir, f)
                    shutil.copy2(full_path, dest)
                    local_files.append(dest)
            return local_files
        else:
            print(f"⚠️ Path is not GS and not a valid local dir: {gcs_path}")
            return []

    # GCS Logic
    gcs_path_clean = gcs_path[5:]
    bucket_name = gcs_path_clean.split('/')[0]
    prefix = '/'.join(gcs_path_clean.split('/')[1:])
    
    storage_client = storage.Client()
    bucket = storage_client.bucket(bucket_name)
    blobs = bucket.list_blobs(prefix=prefix)
    
    downloaded_files = []
    for blob in blobs:
        if blob.name.endswith('/'): continue # skip directories
        
        local_file_path = os.path.join(local_dir, os.path.basename(blob.name))
        blob.download_to_filename(local_file_path)
        downloaded_files.append(local_file_path)
        print(f"✅ Downloaded {blob.name} to {local_file_path}")
        
    return downloaded_files

def upload_to_gcs(local_path: str, bucket_name: str, destination_blob_name: str):
    """Uploads a file to GCS. Falls back to local path if GCS fails."""
    try:
        storage_client = storage.Client()
        bucket = storage_client.bucket(bucket_name)
        blob = bucket.blob(destination_blob_name)
        blob.upload_from_filename(local_path)
        print(f"✅ Uploaded {local_path} to gs://{bucket_name}/{destination_blob_name}")
        return f"gs://{bucket_name}/{destination_blob_name}"
    except Exception as e:
        print(f"🏠 GCS Upload skipped (Local Test Mode): {str(e)}")
        return local_path

def generate_signed_url(bucket_name: str, blob_name: str, expiration_hours: int = 24):
    """Generates a v4 signed URL for downloading a blob. Returns local path if GCS fails."""
    try:
        storage_client = storage.Client()
        bucket = storage_client.bucket(bucket_name)
        blob = bucket.blob(blob_name)

        url = blob.generate_signed_url(
            version="v4",
            expiration=timedelta(hours=expiration_hours),
            method="GET",
        )
        return url
    except Exception as e:
        # If it's a local path from upload_to_gcs, return it as a file URL
        if os.path.exists(blob_name):
            return f"file:///{os.path.abspath(blob_name).replace('\\', '/')}"
        return blob_name

# ==================== PRODUCTION LOGGING (GSheet & Email) ====================

try:
    from googleapiclient.discovery import build
except ImportError:
    build = None

def log_to_google_sheet(job_id, status, details, result_url=""):
    """Logs job status to a central Google Sheet. Fails gracefully if not configured."""
    if not config.LOG_SHEET_ID or not build:
        return
    
    try:
        # Uses default credentials (Workload Identity in Cloud Run)
        service = build('sheets', 'v4')
        range_name = 'Sheet1!A:E'
        values = [[
            datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            job_id,
            status,
            details,
            result_url
        ]]
        body = {'values': values}
        service.spreadsheets().values().append(
            spreadsheetId=config.LOG_SHEET_ID,
            range=range_name,
            valueInputOption='USER_ENTERED',
            body=body
        ).execute()
        print(f"✅ Logged to Google Sheet: {config.LOG_SHEET_ID}")
    except Exception as e:
        print(f"⚠️ Spreadsheet logging skipped/failed: {str(e)}")

def send_email_notification(job_id, status, subject_text, body_text):
    """Placeholder for email notification. Intent only for now."""
    if not config.NOTIFICATION_EMAIL:
        return
    print(f"📧 Notification intended for Job {job_id}: {subject_text}")

# Note: Benchmark data is automatically checked and downloaded inside the processing function for better performance.

DEFAULT_FALLBACK_TAT_RULES = config.DEFAULT_FALLBACK_TAT_RULES

fmt = "%Y-%m-%d %H:%M"

# ==================== EXECUTION LOGGING ====================

# Log directory and consolidated filename
LOG_DIR = config.LOG_DIR

LOG_FILE_PATH = os.path.join(LOG_DIR, "scoping_history.json")

def get_user_email():
    """
    Get the identity of the user.
    In Colab: returns authenticated email.
    On Local: returns OS username.
    """
    try:
        # 1. Try Colab/Google Auth
        import google.auth
        credentials, project = google.auth.default()
        if hasattr(credentials, 'service_account_email'):
            return credentials.service_account_email

        # 2. Try UserInfo endpoint (if token exists)
        try:
            import requests
            token = credentials.token
            response = requests.get(
                'https://www.googleapis.com/oauth2/v1/userinfo',
                headers={'Authorization': f'Bearer {token}'}
            )
            if response.status_code == 200:
                user_info = response.json()
                return user_info.get('email', 'unknown@user.com')
        except:
            pass
    except:
        pass

    # 3. Fallback: Local OS username (Great for Windows/Desktop users)
    try:
        return getpass.getuser()
    except:
        return 'default'

def log_execution(job_ids, document_files, status, error=None, outputs=None, execution_time=None, **kwargs):
    """
    Log execution details to a single consolidated JSON file.
    Appends new entries to a list for easy tracking.
    """
    try:
        user_identity = get_user_email()
        timestamp_str = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        execution_id_str = datetime.now().strftime("%Y%m%d_%H%M%S_%f")

        # Create the new entry
        log_entry = {
            "timestamp": timestamp_str,
            "execution_id": execution_id_str,
            "user_email": user_identity,
            "job_ids": job_ids if isinstance(job_ids, list) else [job_ids],
            "status": status,
            "document_files": [os.path.basename(f) if isinstance(f, str) else f for f in (document_files or [])],
            "num_documents": len(document_files) if document_files else 0,
            "execution_time_seconds": execution_time,
            "error": error,
            "outputs": outputs or {},
            "parameters": kwargs
        }

        # Ensure directory exists
        os.makedirs(LOG_DIR, exist_ok=True)

        # Handle Consolidated JSON File
        all_logs = []
        if os.path.exists(LOG_FILE_PATH):
            try:
                with open(LOG_FILE_PATH, 'r', encoding='utf-8') as f:
                    all_logs = json.load(f)
                    if not isinstance(all_logs, list):
                        all_logs = []
            except (json.JSONDecodeError, Exception):
                # Backup corrupted file if necessary and start fresh
                all_logs = []

        # Append and Save
        all_logs.append(log_entry)
        with open(LOG_FILE_PATH, 'w', encoding='utf-8') as f:
            json.dump(all_logs, f, indent=2, ensure_ascii=False)

        print(f"✅ Execution logged to: {os.path.basename(LOG_FILE_PATH)}")

    except Exception as e:
        print(f"⚠️ Failed to log execution: {e}")

# ==================== PREPROCESSING FUNCTIONS ====================

def analyze_word_document(file_path):
    """Analyze DOCX for preprocessing flags"""
    results = {
        "text": "",
        "comments": [],
        "tracked_changes": [],
        "image_count": 0,
        "ocr_images": [],
        "flags": []
    }
    try:
        doc = docx.Document(file_path)
        results["text"] = "\n".join([p.text for p in doc.paragraphs])
        for section in doc.sections:
            if section.header:
                results["text"] += "\n[Header]\n" + "\n".join([p.text for p in section.header.paragraphs])
            if section.footer:
                results["text"] += "\n[Footer]\n" + "\n".join([p.text for p in section.footer.paragraphs])
    except Exception as e:
        results["flags"].append(f"❌ Could not extract text: {e}")

    for pattern in NON_TRANSLATABLE_PATTERNS:
        if re.search(pattern, results["text"], re.IGNORECASE):
            results["flags"].append("⚠️ Contains 'Do Not Translate' instructions")

    try:
        with zipfile.ZipFile(file_path, 'r') as z:
            ns = {
                "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
                "w15": "http://schemas.microsoft.com/office/word/2012/wordml"
            }
            if "word/comments.xml" in z.namelist():
                comments_xml = z.read("word/comments.xml")
                root = ET.fromstring(comments_xml)
                for c in root.findall(".//w:comment", ns):
                    full_text = "".join([t.text or "" for t in c.findall(".//w:t", ns)])
                    if full_text.strip():
                        results["comments"].append(full_text)
            if "word/commentsExtended.xml" in z.namelist():
                ext_xml = z.read("word/commentsExtended.xml")
                root_ext = ET.fromstring(ext_xml)
                for c in root_ext.findall(".//w15:commentEx", ns):
                    txt = c.get("{http://schemas.microsoft.com/office/word/2012/wordml}text")
                    if txt:
                        results["comments"].append(txt)
            if results["comments"]:
                results["flags"].append("⚠️ Document contains reviewer comments")
    except Exception as e:
        results["flags"].append(f"❌ Error reading comments: {e}")

    try:
        with zipfile.ZipFile(file_path) as z:
            xml = z.read("word/document.xml")
            root = ET.fromstring(xml)
            ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
            ins = root.findall(".//w:ins", ns)
            dels = root.findall(".//w:del", ns)
            if ins or dels:
                results["flags"].append("⚠️ Tracked changes detected")
            for n in ins:
                txt = "".join(t.text or "" for t in n.findall(".//w:t", ns))
                if txt:
                    results["tracked_changes"].append(f"Inserted: {txt}")
            for n in dels:
                txt = "".join(t.text or "" for t in n.findall(".//w:t", ns))
                if txt:
                    results["tracked_changes"].append(f"Deleted: {txt}")
    except Exception as e:
        results["flags"].append(f"❌ Could not check tracked changes: {e}")

    try:
        with zipfile.ZipFile(file_path) as z:
            imgs = [f for f in z.namelist() if f.startswith("word/media/")]
            results["image_count"] = len(imgs)
            if imgs:
                results["flags"].append(f"⚠️ Contains {results['image_count']} images (may include text)")
            for img_name in imgs:
                img_bytes = z.read(img_name)
                try:
                    img = Image.open(io.BytesIO(img_bytes))
                    ocr_text = pytesseract.image_to_string(img).strip()
                    has_text = len(ocr_text) > 0
                    lang = "unknown"
                    if has_text:
                        try:
                            lang = detect(ocr_text)
                        except LangDetectException:
                            lang = "unknown"
                        results["flags"].append(f"⚠️ Image {img_name} contains text (lang: {lang})")
                    results["ocr_images"].append({
                        "image_name": img_name.split("/")[-1],
                        "has_text": has_text,
                        "language": lang
                    })
                except Exception:
                    results["ocr_images"].append({
                        "image_name": img_name.split("/")[-1],
                        "has_text": False,
                        "language": "error"
                    })
    except Exception as e:
        results["flags"].append(f"❌ Error processing images: {e}")

    if results["flags"]:
        results["flags"].append("✅ Manual preprocessing recommended")
    return results

def analyze_pptx(file_path):
    """Analyze PPTX for preprocessing flags"""
    results = {
        "slides": [],
        "comments": [],
        "image_count": 0,
        "image_text_flags": [],
        "flags": []
    }
    try:
        with zipfile.ZipFile(file_path, 'r') as z:
            slide_files = sorted([f for f in z.namelist() if f.startswith("ppt/slides/slide") and f.endswith(".xml")])
            for idx, slide_file in enumerate(slide_files, start=1):
                slide_text = ""
                notes_text = ""
                xml_content = z.read(slide_file)
                root = ET.fromstring(xml_content)
                ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                      "p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
                slide_texts = [t.text for t in root.findall(".//a:t", ns) if t.text]
                slide_text += "\n".join(slide_texts)
                notes_file = f"ppt/notesSlides/notesSlide{idx}.xml"
                if notes_file in z.namelist():
                    notes_xml = z.read(notes_file)
                    root_notes = ET.fromstring(notes_xml)
                    note_texts = [t.text for t in root_notes.findall(".//a:t", ns) if t.text]
                    notes_text = "\n".join(note_texts)
                    if notes_text:
                        slide_text += "\n[Notes]\n" + notes_text
                results["slides"].append(slide_text)
                for pattern in NON_TRANSLATABLE_PATTERNS:
                    if re.search(pattern, slide_text, re.IGNORECASE):
                        if notes_text:
                            results["flags"].append(f"⚠️ Notes on Slide {idx} contain non-translatable instructions")
                        else:
                            results["flags"].append(f"⚠️ Slide {idx} contains non-translatable instructions")
            comment_files = sorted([f for f in z.namelist() if f.startswith("ppt/comments/comment") and f.endswith(".xml")])
            for cfile in comment_files:
                xml_content = z.read(cfile)
                root = ET.fromstring(xml_content)
                ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
                for comment in root.findall(".//p:cm", ns):
                    text = comment.attrib.get("text", "")
                    author = comment.attrib.get("author", "Unknown")
                    slide_idx = int(comment.attrib.get("parentSlide", "0"))
                    results["comments"].append((slide_idx, text, author))
            if results["comments"]:
                results["flags"].append("⚠️ Contains comments")
            img_files = [f for f in z.namelist() if f.startswith("ppt/media/")]
            results["image_count"] = len(img_files)
            if img_files:
                results["flags"].append(f"⚠️ Contains {len(img_files)} images")
            for img_file in img_files:
                img_data = z.read(img_file)
                try:
                    img = Image.open(io.BytesIO(img_data))
                    ocr_text = pytesseract.image_to_string(img).strip()
                    has_text = bool(ocr_text)
                    lang = "unknown"
                    if has_text:
                        try:
                            lang = detect(ocr_text)
                        except:
                            lang = "unknown"
                        results["flags"].append(f"⚠️ Image {img_file} contains text (lang: {lang})")
                    results["image_text_flags"].append((img_file, has_text, lang))
                except Exception as e:
                    results["flags"].append(f"❌ Could not OCR {img_file}: {e}")
                    results["image_text_flags"].append((img_file, False, "unknown"))
    except Exception as e:
        results["flags"].append(f"❌ Could not process PPTX: {e}")
    if results["flags"]:
        results["flags"].append("✅ Manual preprocessing recommended")
    return results

def analyze_excel(file_path):
    """Analyze Excel for preprocessing flags"""
    results = {
        "sheets": [],
        "flags": [],
        "image_count": 0,
        "ocr_images": [],
        "cell_comments": []
    }
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        results["sheets"] = wb.sheetnames
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=False):
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        for pattern in NON_TRANSLATABLE_PATTERNS:
                            if re.search(pattern, cell.value, re.IGNORECASE):
                                results["flags"].append(f"⚠️ Sheet '{sheet}' cell {cell.coordinate} contains non-translatable instructions")
                    if cell.comment:
                        results["cell_comments"].append({
                            "sheet": sheet,
                            "cell": cell.coordinate,
                            "author": cell.comment.author,
                            "text": cell.comment.text
                        })
                        results["flags"].append(f"⚠️ Sheet '{sheet}' cell {cell.coordinate} has a comment by {cell.comment.author}")
        with zipfile.ZipFile(file_path, 'r') as z:
            img_files = [f for f in z.namelist() if f.startswith("xl/media/")]
            results["image_count"] = len(img_files)
            if img_files:
                results["flags"].append(f"⚠️ Contains {len(img_files)} images (may include text)")
            for img_file in img_files:
                img_bytes = z.read(img_file)
                try:
                    img = Image.open(io.BytesIO(img_bytes))
                    ocr_text = pytesseract.image_to_string(img).strip()
                    has_text = len(ocr_text) > 0
                    lang = "unknown"
                    preview = ""
                    if has_text:
                        try:
                            lang = detect(ocr_text)
                        except LangDetectException:
                            lang = "unknown"
                        preview = " ".join(ocr_text.split())[:50]
                        results["flags"].append(f"⚠️ Excel image {img_file} contains text (lang: {lang})")
                    results["ocr_images"].append({
                        "image_name": img_file.split("/")[-1],
                        "has_text": has_text,
                        "language": lang,
                        "text_preview": preview
                    })
                except:
                    results["ocr_images"].append({
                        "image_name": img_file.split("/")[-1],
                        "has_text": False,
                        "language": "error",
                        "text_preview": ""
                    })
    except Exception as e:
        results["flags"].append(f"❌ Could not process Excel: {e}")
    if results["flags"]:
        results["flags"].append("✅ Manual preprocessing recommended")
    return results

def preprocess_file(file_path):
    """Wrapper for preprocessing checks"""
    filename_lower = file_path.lower()

    # Automatically flag INDD/IDML files for manual preprocessing
    if filename_lower.endswith(".indd") or filename_lower.endswith(".idml"):
        return {
            "flags": [
                "⚠️ INDD/IDML file format detected",
                "⚠️ This file type requires manual preprocessing",
                "✅ Manual preprocessing recommended"
            ]
        }

    if filename_lower.endswith(".docx"):
        return analyze_word_document(file_path)
    elif filename_lower.endswith(".pptx"):
        return analyze_pptx(file_path)
    elif filename_lower.endswith(".xlsx"):
        return analyze_excel(file_path)
    else:
        return {"flags": []}

def get_preprocessing_report(filename, result):
    """Generate preprocessing report"""
    lines = []
    lines.append(f"\n📂 Preprocessing: {filename}")
    lines.append("\n============== ANALYSIS REPORT ==============")

    # 1. Specialized Reports
    if filename.lower().endswith(".indd") or filename.lower().endswith(".idml"):
        lines.append("\n⚠️ File Type: INDD/IDML")
        lines.append("⚠️ Status: Requires Manual Preprocessing")
        lines.append("\nThis file format cannot be automatically analyzed.")
        lines.append("Please process manually before translation.")
    elif filename.lower().endswith(".docx"):
        lines.append(f"\n📝 Comments found: {len(result.get('comments', []))}")
        lines.append(f"🔄 Tracked changes found: {len(result.get('tracked_changes', []))}")
        lines.append(f"🖼️ Image count: {result.get('image_count', 0)}")
        if result.get("ocr_images"):
            lines.append("\n📸 OCR Image Results:")
            for i in result["ocr_images"]:
                lines.append(str(i))
    elif filename.lower().endswith(".pptx"):
        lines.append(f"\n📝 Total slides: {len(result.get('slides', []))}")
        lines.append(f"📝 Total comments: {len(result.get('comments', []))}")
        lines.append(f"🖼️ Total images: {result.get('image_count', 0)}\n")
        if result.get("image_text_flags"):
            lines.append("🖼️ Images with OCR detection:")
            for img_file, has_text, lang in result["image_text_flags"]:
                if has_text:
                    lines.append(f"- {img_file}: has_text={has_text}, lang={lang}")
    elif filename.lower().endswith(".xlsx"):
        lines.append(f"\n📝 Sheets: {', '.join(result.get('sheets', []))}")
        lines.append(f"📝 Cell comments found: {len(result.get('cell_comments', []))}")
        lines.append(f"🖼️ Images: {result.get('image_count', 0)}")
        if result.get("ocr_images"):
            lines.append("📸 OCR Image Results:")
            for img in result["ocr_images"]:
                if img['has_text']:
                    preview = img['text_preview'].replace("\n", " ")[:80]
                    lines.append(f"- {img['image_name']}: lang={img['language']}, text_preview={preview}...")

    # 2. Add ANY generic flags
    other_flags = [f for f in result.get("flags", []) if "Manual preprocessing recommended" not in f]
    if other_flags:
        lines.append("\n🚩 Flags Detected:")
        for flag in other_flags:
            lines.append(f"- {flag}")

    lines.append("=============================================")
    
    # 3. Success Message
    if not result.get("flags"):
        lines.append("\n✅ No issues found. This document is clean.")
        
    return "\n".join(lines)

# ==================== HELPER FUNCTIONS ====================

def extract_text_multi(file_bytes, filename):
    """Extract text from various file formats"""
    filename = filename.lower()
    if filename.endswith('.pdf'):
        with pdfplumber.open(file_bytes) as pdf:
            return "\n".join(page.extract_text() or '' for page in pdf.pages)
    elif filename.endswith('.docx'):
        if hasattr(file_bytes, "seek"):
            file_bytes.seek(0)
        doc = docx.Document(file_bytes)
        texts = [para.text for para in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    texts.append(cell.text)
        return "\n".join(filter(None, texts))
    elif filename.endswith('.xml'):
        if hasattr(file_bytes, "seek"):
            file_bytes.seek(0)
        tree = ET.parse(file_bytes)
        root = tree.getroot()
        def local_name(tag):
            if tag.startswith('{'):
                return tag.split('}', 1)[1]
            return tag
        texts = []
        if local_name(root.tag).lower() == "tms":
            for node in root.findall('.//tmsnode'):
                chunk = " ".join(t.strip() for t in node.itertext() if t.strip())
                if chunk:
                    texts.append(chunk)
        if not texts:
            texts = [" ".join(t.strip() for t in root.itertext() if t.strip())]
        return "\n".join(filter(None, texts))
    elif filename.endswith('.pptx'):
        prs = Presentation(file_bytes)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    elif filename.endswith('.txt'):
        return file_bytes.read().decode('utf-8')
    elif filename.endswith('.json+html') or filename.endswith('.json'):
        raw = json.load(file_bytes)
        title = raw.get("title", "")
        body = raw.get("body", "")
        soup = BeautifulSoup(body, "html.parser")
        return title + "\n" + soup.get_text(separator="\n")
    elif filename.endswith('.xlsx'):
        xls = pd.ExcelFile(file_bytes)
        all_text = ""
        for sheet in xls.sheet_names:
            df_sheet = pd.read_excel(xls, sheet_name=sheet)
            sheet_text = " ".join(df_sheet.astype(str).values.flatten())
            all_text += sheet_text + "\n"
        return all_text
    elif filename.endswith('.idml'):
        all_text = ""
        with zipfile.ZipFile(file_bytes) as z:
            for fname in z.namelist():
                if fname.endswith('.xml'):
                    with z.open(fname) as f:
                        tree = ET.parse(f)
                        root = tree.getroot()
                        for story in root.findall('.//Story'):
                            all_text += "".join(story.itertext()) + "\n"
        return all_text
    elif filename.endswith('.xliff') or filename.endswith('.pptx.xliff') or filename.endswith('.sdlxliff'):
        all_text = ""
        tree = ET.parse(file_bytes)
        root = tree.getroot()
        ns = {}
        if root.tag.startswith('{') and '}' in root.tag:
            uri = root.tag[root.tag.find('{') + 1: root.tag.find('}')]
            ns = {'ns': uri}
            source_xpath = './/ns:source'
            target_xpath = './/ns:target'
        else:
            source_xpath = './/source'
            target_xpath = './/target'
        for elem in root.findall(source_xpath, ns):
            text = ''.join(elem.itertext()).strip()
            if text:
                all_text += text + "\n"
        for elem in root.findall(target_xpath, ns):
            text = ''.join(elem.itertext()).strip()
            if text:
                all_text += text + "\n"
        if all_text.strip():
            return all_text
        internal_xpath = './/ns:internal-file' if ns else './/internal-file'
        for internal in root.findall(internal_xpath, ns):
            form = internal.attrib.get('form', '').lower()
            if form == 'base64':
                data = ''.join(internal.itertext()).strip()
                if not data:
                    continue
                try:
                    decoded = base64.b64decode(data)
                    embedded_name = internal.attrib.get('original', 'embedded.docx')
                    embedded_stream = io.BytesIO(decoded)
                    embedded_stream.seek(0)
                    return extract_text_multi(embedded_stream, embedded_name)
                except Exception as exc:
                    print(f"[WARNING] Failed to decode embedded file in {filename}: {exc}")
        return all_text
    elif filename.endswith('.liltjson'):
        raw = json.load(file_bytes)
        all_text = ""
        if isinstance(raw, dict) and "document" in raw:
            for item in raw["document"]:
                val = item.get("value", "")
                if isinstance(val, str):
                    soup = BeautifulSoup(val, "html.parser")
                    text = soup.get_text(separator="\n")
                    all_text += text + "\n"
        return all_text
    elif filename.endswith('.srt'):
        text = ""
        for line in file_bytes.read().decode('utf-8').splitlines():
            if line.strip() == "" or line.strip().isdigit() or "-->" in line:
                continue
            text += line + "\n"
        return text
    else:
        raise ValueError(f"Unsupported file type: {filename}")

EXCEL_UNSAFE_PREFIXES = ("=", "+", "-", "@")

def sanitize_for_excel(text):
    """Prevent Excel from treating text as a formula"""
    if text is None:
        return ""
    text = str(text).replace("\r", "")
    stripped = text.lstrip()
    if stripped.startswith(EXCEL_UNSAFE_PREFIXES):
        return "'" + text
    return text

def load_fallback_sla_rules(path=FALLBACK_SLA_PATH):
    """Load fallback SLA rules from JSON"""
    try:
        if os.path.exists(path):
            with open(path, "r", encoding="utf-8") as f:
                data = json.load(f)
            return data.get("project_due_date_list", DEFAULT_FALLBACK_TAT_RULES)
    except Exception as exc:
        print(f"[WARNING] Failed to load fallback SLA file: {exc}")
    return DEFAULT_FALLBACK_TAT_RULES

def select_json_tat(total_words, fallback_rules):
    """Select TAT hours from fallback rules based on word volume"""
    for rule in fallback_rules:
        min_words = rule.get("wordVolumeMin", 0)
        max_words = rule.get("wordVolumeMax", -1)
        max_ok = max_words == -1 or total_words <= max_words
        if total_words >= min_words and max_ok:
            return float(rule.get("hoursUntilDue", 96))
    return float(DEFAULT_FALLBACK_TAT_RULES[-1]["hoursUntilDue"])

def compute_ramped_tat(total_words, daily_throughput, ramp_days):
    """Compute TAT hours using ramped throughput logic"""
    if daily_throughput is None or ramp_days is None:
        return None
    daily_throughput = float(daily_throughput)
    ramp_days = max(1, int(ramp_days))
    if daily_throughput <= 0:
        raise ValueError("Daily throughput must be > 0.")
    ramp_rate = daily_throughput / ramp_days
    words_remaining = float(total_words)
    hours = 0.0
    for _ in range(ramp_days):
        if words_remaining <= 0:
            break
        produced = min(ramp_rate, words_remaining)
        hours += (produced / ramp_rate) * 24
        words_remaining -= produced
    if words_remaining > 0:
        hours += (words_remaining / daily_throughput) * 24
    return max(1, math.ceil(hours))

def format_tat(hours):
    """Format TAT hours to readable format"""
    days = hours / 24
    if days.is_integer():
        days_str = f"{int(days)} day"
        if days != 1:
            days_str += "s"
    else:
        days_str = f"{days:.2f} days"
    return f"{int(hours)} hrs ({days_str})"

def add_business_hours(start, hours):
    """Add business hours skipping weekends"""
    current = start
    remaining_hours = hours
    while remaining_hours > 0:
        while current.weekday() >= 5:
            current += timedelta(days=1)
            current = current.replace(hour=current.hour, minute=current.minute, second=current.second, microsecond=0)
        current += timedelta(hours=1)
        remaining_hours -= 1
    return current

def compute_sla_tat(total_words, sla_min_volume, sla_max_volume, sla_tat_in_hours):
    """Compute SLA TAT with fallback logic"""
    try:
        v_min = float(sla_min_volume) if sla_min_volume is not None else None
    except:
        v_min = None
    try:
        v_max = float(sla_max_volume) if sla_max_volume is not None else None
    except:
        v_max = None
    if v_min is not None and v_max is not None:
        if v_min <= total_words <= v_max:
            return float(sla_tat_in_hours)
    if total_words > 10000:
        blocks = math.ceil((total_words - 10000) / 5000)
        tat = 96 + blocks * 24
        print(f"[FALLBACK] {total_words} words → {tat} hours")
        return tat
    print(f"[WARNING] No SLA matched for {total_words} words (range={v_min}-{v_max}). Using provided SLA value.")
    return float(sla_tat_in_hours)





# ==================== MAIN PROCESSING FUNCTION ====================

def process_translation_project(request: ScopingRequest):
    """
    Orchestrates the scoping analysis backend.
    Takes a ScopingRequest and processes files from GCS.
    """
    start_time = datetime.now()
    timestamp_str = start_time.strftime("%Y%m%d_%H%M%S")
    job_dir = os.path.join(DATA_BASE_DIR, f"job_{timestamp_str}")
    os.makedirs(job_dir, exist_ok=True)
    
    input_dir = os.path.join(job_dir, "inputs")
    output_dir = os.path.join(job_dir, "outputs")
    os.makedirs(input_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    try:
        status = "🔧 Setting up environment...\n"
        api_key = config.OPENAI_API_KEY
        if not api_key:
            return "❌ Error: OpenAI API Key is missing.", None, None, None

        # Unpack request
        job_ids_input = request.job_ids
        gcs_input_path = request.gcs_input_path
        user_instructions = request.instructions
        ramped_daily_throughput = request.ramped_daily_throughput
        ramp_up_days = request.ramp_up_days
        translation_pct = request.translator_pct
        review_pct = request.reviewer_pct
        pm_pct = request.pm_pct

        # ==================== DOWNLOAD INPUTS FROM GCS ====================
        status += f"📥 Downloading documents from {gcs_input_path}...\n"
        document_files = download_from_gcs(gcs_input_path, input_dir)
        if not document_files:
            return f"❌ Error: No files found at {gcs_input_path}", None, None, None

        status += f"✅ Downloaded {len(document_files)} documents\n"

        # Log start to Google Sheet
        log_to_google_sheet(timestamp_str, "IN_PROGRESS", f"Processing {len(document_files)} documents for Job IDs: {job_ids_input}")

        user_ramp_config = None
        if ramped_daily_throughput and ramp_up_days:
            try:
                throughput_val = float(ramped_daily_throughput)
                ramp_days_val = int(ramp_up_days)
                if throughput_val > 0 and ramp_days_val > 0:
                    user_ramp_config = {
                        "throughput": throughput_val,
                        "ramp_days": ramp_days_val
                    }
                    status += f"✅ User throughput override: {throughput_val} words/day after {ramp_days_val} day(s) ramp\n"
                else:
                    status += "⚠️ Invalid ramp inputs provided — ignoring user override\n"
            except Exception:
                status += "⚠️ Could not parse ramp inputs — ignoring user override\n"
        else:
            status += "ℹ️ No ramped throughput override provided\n"

        job_ids = [int(jid.strip()) for jid in job_ids_input.split(',') if jid.strip().isdigit()]
        if not job_ids:
            return "❌ Error: No valid job IDs provided", None, None, None

        user_project_input = ",".join(map(str, job_ids))
        status += f"✅ Job IDs: {user_project_input}\n"

        # ==================== PREPROCESSING PHASE ====================
        status += "\n📋 Running preprocessing checks...\n"
        preprocessing_report = ""
        preprocessing_flags_found = False

        if document_files:
            for doc_file_path in document_files:
                filename = os.path.basename(doc_file_path)
                preprocess_result = preprocess_file(doc_file_path)

                if preprocess_result.get("flags"):
                    preprocessing_flags_found = True

                preprocessing_report += get_preprocessing_report(filename, preprocess_result)
                preprocessing_report += "\n"

        if preprocessing_flags_found:
            status += "⚠️ Preprocessing flags detected - review recommended\n"
        else:
            status += "✅ No preprocessing issues found\n"

        status += "\n" + preprocessing_report + "\n"

        # ==================== LOAD BENCHMARK DATA ====================
        status += "📊 Loading benchmark data...\n"

        try:
            # Check if file exists and is valid; download if necessary
            if not is_valid_parquet(BENCHMARK_LOCAL_PATH):
                status += "🔄 Benchmark file missing or corrupted. Downloading from source...\n"
                if os.path.exists(BENCHMARK_LOCAL_PATH):
                    os.remove(BENCHMARK_LOCAL_PATH)
                download_file_from_google_drive(BENCHMARK_FILE_ID, BENCHMARK_LOCAL_PATH)

            # Use local path unless an override is provided via UI
            effective_benchmark_path = BENCHMARK_LOCAL_PATH
            if benchmark_file is not None:
                effective_benchmark_path = getattr(benchmark_file, 'name', benchmark_file)

            df = pd.read_parquet(effective_benchmark_path)
            status += f"✅ Loaded {len(df)} benchmark records\n"
        except Exception as e:
            return f"❌ Error loading benchmark: {str(e)}", None, None, None
        fallback_rules = load_fallback_sla_rules()
        status += f"✅ Loaded fallback SLA rules ({len(fallback_rules)} buckets)\n"

        # ==================== QUERY BIGQUERY ====================
        status += "🔍 Querying BigQuery...\n"
        
        try:
            # Explicitly check for service account credentials
            if config.GOOGLE_APPLICATION_CREDENTIALS and os.path.exists(config.GOOGLE_APPLICATION_CREDENTIALS):
                bq_client = bigquery.Client.from_service_account_json(config.GOOGLE_APPLICATION_CREDENTIALS, project=config.PROJECT_ID)
            else:
                bq_client = bigquery.Client(project=config.PROJECT_ID)
        except Exception as e:
            if "DefaultCredentialsError" in str(e):
                return "❌ BigQuery Auth Error: No credentials found. Docker is isolated, so it needs a Service Account JSON key to talk to BigQuery. Please see README.md for the 'docker run' command with keys.", None, None, None
            return f"❌ BigQuery Error: {str(e)}", None, None, None

        assignment_sql = f"""
        WITH lang_agg AS (
          SELECT
            CONCAT(
              m.trgLang,
              IF(m.trgLocale IS NOT NULL AND m.trgLocale != '', CONCAT('-', m.trgLocale), '')
            ) AS target_lang_full,
            SUM(
              IFNULL(m.onlineFuzzy75to84Words,0) +
              IFNULL(m.onlineFuzzy85to94Words,0) +
              IFNULL(m.onlineFuzzy95to99Words,0) +
              IFNULL(m.onlineNewWords,0)
            ) AS total_translated_words,
            MIN(m.projectCreatedDate) AS projectCreatedDate,
            MAX(m.dueDate) AS dueDate,
            ARRAY_AGG(DISTINCT m.jobId) AS jobId,
            ARRAY_AGG(DISTINCT m.ProjectID) AS ProjectID,
            ARRAY_AGG(DISTINCT m.projectName) AS projectName,
            ARRAY_AGG(DISTINCT m.srcLang) AS srcLang,
            ARRAY_AGG(DISTINCT m.workflow) AS workflow,
            ARRAY_AGG(DISTINCT IFNULL(m.domain_name,'UNKNOWN')) AS domain_name,
            ANY_VALUE(m.salesforceCustomerId) AS salesforceCustomerId,
            ANY_VALUE(m.CustomerName) AS customer_name
          FROM `Datastudio.L1toL3ProductionMetrics` m
          WHERE jobId IN ({user_project_input})
          GROUP BY target_lang_full
        )
        SELECT
          la.* EXCEPT(salesforceCustomerId),
          sla.tat_in_hours__c,
          sla.min_volume__c,
          sla.max_volume__c
        FROM lang_agg la
        LEFT JOIN `salesforce.Lilt_Org__c` o
          ON la.salesforceCustomerId = o.account__c
            AND (o.status__c = 'Active' OR o.status__c IS NULL)
        LEFT JOIN `salesforce.Service_Level_Agreement__c` sla
          ON o.id = sla.lilt_org__c
            AND la.total_translated_words >= sla.min_volume__c
            AND (la.total_translated_words <= sla.max_volume__c OR sla.max_volume__c IS NULL)
            AND (sla.status__c = 'Active' OR sla.status__c IS NULL)
        ORDER BY total_translated_words DESC;
        """
        df_assignment = bq_client.query(assignment_sql).to_dataframe()
        status += f"✅ Fetched {len(df_assignment)} project assignments\n"

        # ==================== EXTRACT DOCUMENT TEXT ====================
        status += "📄 Extracting text from documents...\n"
        if not document_files:
            return "❌ Error: At least one document file is required", None, None, None

        all_texts = []
        total_word_count = 0
        for doc_file_path in document_files:
            with open(doc_file_path, "rb") as f:
                text = extract_text_multi(f, doc_file_path)
                all_texts.append(text)
                total_word_count += len(text.split())

        merged_text = "\n".join(all_texts)
        doc_text = "\n".join([line.strip() for line in merged_text.splitlines() if line.strip()])
        status += f"✅ Extracted {total_word_count} words from {len(document_files)} file(s)\n"

        # ==================== GPT-4O ANALYSIS ====================
        status += "🤖 Analyzing document with GPT-4o...\n"
        valid_domains = df['domain'].dropna().unique().tolist()
        valid_content_types = df['content_type'].dropna().unique().tolist()


        system_prompt = f"""
        {{
          "content_type": "...",        # Must be one of these: {valid_content_types}
          "domain": "...",              # Must be one of these: {valid_domains}
          "purpose": "...",
          "key_terms": [...],
          "complexity": {{

              "sentence_structure": "Easy/Medium/Hard",
              "word_difficulty": "Easy/Medium/Hard",
              "tone_specificity": "...",
              "idioms_present": true/false,
              "double_meanings": true/false,
              "formatting_or_tags": true/false,
              "special_instructions_applied": true/false
          }},
          "quality_considerations": []
        }}

        Instructions:

        1. Evaluate complexity exactly as described in the "complexity" field.
          - Consider sentence length, word familiarity, tone, idioms, double meanings, and formatting indicators.

        2. Generate **dynamic, actionable notes** in "quality_considerations" based on document content.

        3. Based on the document characteristics and complexity, generate a single unified sourcing criteria section. Do not split it into Translator or Reviewer roles. Instead, provide one consolidated set of requirements that combines all relevant linguistic expectations.

          The unified sourcing criteria must cover:

          - Required linguistic expertise or domain knowledge (e.g., tech, marketing, legal)
          - Preferred experience level (e.g., high proficiency, prior domain-specific experience)
          - Tone or style familiarity needed (casual, instructional, brand-specific)
          - Sensitivity to idioms, cultural context, regional variants, and double meanings
          - Familiarity with formatting rules, terminology usage, or CAT-tool expectations
          - Ability to ensure accuracy, consistency, tone validation, and terminology control
          - Strengths related to quality checking, clarity, stylistic alignment, and adherence to brand or platform guidelines

            The entire sourcing criteria response must be returned strictly as a **single bullet-point list**, forming one unified block of guidance suitable for any linguist involved in the document workflow.

        5. Always choose "content_type" and "domain" strictly from the provided lists.

        6. Return ONLY a JSON object, with **only the applicable keys**:

        """
        example_json = """
        {
          "general_sourcing_criteria": "- bullet A\\n- bullet B\\n- bullet C"
        }
        """
        final_prompt = system_prompt + example_json

        client = openai.OpenAI(api_key=api_key)
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": final_prompt},
                {"role": "user", "content": doc_text[:12000]}
            ],
            temperature=0
        )


        model_response = response.choices[0].message.content
        raw_output = model_response.strip()
        status += "✅ GPT-4o analysis complete\n"

        # ==================== PARSE GPT RESPONSE ====================
        status += "🔍 Parsing analysis results...\n"
        cleaned_output = raw_output.strip()
        cleaned_output = re.sub(r"^```json", "", cleaned_output)
        cleaned_output = re.sub(r"```$", "", cleaned_output)
        cleaned_output = cleaned_output.strip()
        cleaned_output = re.sub(
            r'(?<=: )"(.*?)"',
            lambda m: '"' + m.group(1).replace('\n', '\\n').replace('\r', '') + '"',
            cleaned_output,
            flags=re.DOTALL
        )

        try:
            analysis_json = json.loads(cleaned_output)
        except json.JSONDecodeError:
            try:
                analysis_json = ast.literal_eval(cleaned_output)
            except Exception as e:
                return f"❌ Error parsing GPT response: {str(e)}\n\nRaw output:\n{cleaned_output[:1000]}", None, None, None

        gpt_domain_check = str(analysis_json.get("domain", "") or "").strip()
        gpt_ct_check = str(analysis_json.get("content_type", "") or "").strip()
        if not gpt_domain_check or not gpt_ct_check:
            status += "⚠️ GPT did not reliably return 'domain' and/or 'content_type'.\n"

        general_assignee_instructions = sanitize_for_excel(
            analysis_json.get("general_sourcing_criteria", "N/A")
        )

        auto_Instructions = "; ".join(analysis_json.get("quality_considerations", []))
        combined_instructions = auto_Instructions
        if user_instructions:
            combined_instructions += f"; USER NOTE: {user_instructions}"
        combined_instructions = sanitize_for_excel(combined_instructions or "")

        complexity = analysis_json.get("complexity", {})
        complexity_str = json.dumps(complexity)

        filtered_df_global = df
        ref_customers = []
        mean_wph = 0
        multiplier = 1

        productivity_json = {
            "reference_customers": ref_customers,
            "estimated_words_per_hour": mean_wph,
            "documents_multiplier_lilt_vs_industry_standard": multiplier
        }

        status += "📋 Building project summaries...\n"
        summary_rows = []

        for _, row in tqdm(df_assignment.iterrows(), total=len(df_assignment), desc="Processing Workflows"):
            gpt_domain = analysis_json.get("domain", "").strip().lower()
            bq_domains = [str(d).strip().lower() for d in row.get("domain_name", ["UNKNOWN"])]

            row_domains = []
            for d in bq_domains:
                if d != "unknown" and d != "":
                    row_domains.append(d)

            if not row_domains and gpt_domain:
                row_domains = [gpt_domain]

            if not row_domains:
                row_domains = ["unknown"]

            content_type = analysis_json.get("content_type", "").strip().lower()

            filtered_df = df[
                (df['domain'].fillna("").str.lower().isin(row_domains)) &
                (df['content_type'].fillna("").str.lower() == content_type)
            ]

            if filtered_df.empty:
                filtered_df = df

            final_domain = ", ".join(row_domains)

            job_id = ", ".join(str(i) for i in row["jobId"])
            ProjectID = ", ".join(str(i) for i in row["ProjectID"])
            srcLang = ", ".join(str(i) for i in row["srcLang"])
            projectName = "\n- ".join(str(i) for i in row["projectName"])
            workflow = ", ".join(str(i) for i in row["workflow"])
            target_lang_full = row["target_lang_full"]
            total_words = int(row["total_translated_words"])
            projectCreatedDate = pd.to_datetime(row["projectCreatedDate"], dayfirst=True)
            actual_due = pd.to_datetime(row["dueDate"], dayfirst=True, errors='coerce')
            query_exec_dt = datetime.now()
            sla_tat_numeric = float(row["tat_in_hours__c"] or 0)
            sla_min_volume = row["min_volume__c"]
            sla_max_volume = row["max_volume__c"]

            errors = []
            if sla_min_volume is None or sla_min_volume < 0:
                errors.append(f"Invalid minimum volume: {sla_min_volume} (row jobId={row['jobId']})")
            if sla_max_volume is not None and sla_max_volume < 0:
                errors.append(f"Invalid maximum volume: {sla_max_volume} (row jobId={row['jobId']})")

            if errors:
                raise ValueError(" ; ".join(errors))

            final_tat_hours = None
            tat_source = "SFDC SLA"

            if user_ramp_config:
                final_tat_hours = compute_ramped_tat(
                    total_words,
                    user_ramp_config["throughput"],
                    user_ramp_config["ramp_days"]
                )
                tat_source = "User Input"
            elif sla_tat_numeric and sla_tat_numeric > 0:
                final_tat_hours = compute_sla_tat(
                    total_words,
                    sla_min_volume,
                    sla_max_volume,
                    sla_tat_numeric
                )

            if final_tat_hours is None:
                final_tat_hours = select_json_tat(total_words, fallback_rules)
                tat_source = "JSON Fallback"

            customer_name = row["customer_name"]
            suggested_due_raw = add_business_hours(projectCreatedDate, final_tat_hours)

            # ==================== SLA SPLIT LOGIC ====================

            # Helper: Validates user input. Returns 0.0 if None/Invalid, unless ALL are None.
            def _safe_float(val):
                try:
                    return float(val) if val is not None and str(val).strip() != "" else None
                except:
                    return None

            t_val = _safe_float(translation_pct_gradio)
            r_val = _safe_float(review_pct_gradio)
            p_val = _safe_float(pm_pct_gradio)

            # If ALL are None/Empty, then we use auto-logic.
            # If AT LEAST ONE is provided, we treat missing ones as 0.0 and attempt to use manual override.
            if t_val is None and r_val is None and p_val is None:
                use_defaults = True
                status += "ℹ️ No manual time splits provided. Using auto-logic.\n"
            else:
                use_defaults = False
                translation_pct = t_val if t_val is not None else 0.0
                review_pct = r_val if r_val is not None else 0.0
                pm_pct = p_val if p_val is not None else 0.0
                status += f"ℹ️ Using manual time split: T:{translation_pct:.0%}, R:{review_pct:.0%}, PM:{pm_pct:.0%}\n"

            if use_defaults:
                workflow_lower = workflow.lower()
                # --- NEW LOGIC: Always 15% PM Time ---
                # CRITICAL: Check specific "Customer Review" workflows BEFORE generic "Translate > Review"
                if "translate > customer review" in workflow_lower or "secondary review" in workflow_lower:
                    # Translate (85%), Passthrough Review (0%), PM (15%)
                    translation_pct, review_pct, pm_pct = 0.85, 0.0, 0.15
                    status += f"ℹ️ Applying Translation-Only split for '{workflow}' (T:85%, R:0%, PM:15%).\n"
                elif "translate > review" in workflow_lower or "prompt response > prompt review" in workflow_lower:
                    # Translate (60%), Review (25%), PM (15%)
                    translation_pct, review_pct, pm_pct = 0.60, 0.25, 0.15
                    status += f"ℹ️ Applying Standard split for '{workflow}' (T:60%, R:25%, PM:15%).\n"
                elif "ai > review" in workflow_lower or "instant review" in workflow_lower or "source review" in workflow_lower:
                    # Passthrough Translate (0%), Review (85%), PM (15%)
                    translation_pct, review_pct, pm_pct = 0.0, 0.85, 0.15
                    status += f"ℹ️ Applying Review-Only split for '{workflow}' (T:0%, R:85%, PM:15%).\n"
                else:
                    # Default fallback: Standard
                    translation_pct, review_pct, pm_pct = 0.60, 0.25, 0.15
                    status += "ℹ️ No specific workflow match, using Standard 60/25/15 split.\n"

            # Normalize to ensure sum is 1.0 (just in case of floating point drift)
            total_pct = translation_pct + review_pct + pm_pct
            if not math.isclose(total_pct, 1.0, rel_tol=1e-9):
                if total_pct > 0:
                    translation_pct /= total_pct
                    review_pct /= total_pct
                    pm_pct /= total_pct
                else:
                    translation_pct, review_pct, pm_pct = 0.60, 0.25, 0.15
                    status += "⚠️ Warning: All user-provided percentages zero. Falling back to default 60/25/15.\n"

            translation_hours = final_tat_hours * translation_pct
            review_hours = final_tat_hours * review_pct
            pm_hours = final_tat_hours * pm_pct

            # Calculate DUE DATES (Datetime objects)
            translation_due_from_creation = add_business_hours(projectCreatedDate, translation_hours)
            review_due_from_creation = add_business_hours(translation_due_from_creation, review_hours)
            pm_due_from_creation = add_business_hours(review_due_from_creation, pm_hours)

            suggested_due = suggested_due_raw.strftime(fmt)
            effective_due = add_business_hours(datetime.now(), final_tat_hours)
            effective_due_str = effective_due.strftime(fmt)

            translation_due_from_execution = add_business_hours(query_exec_dt, translation_hours)
            review_due_from_execution = add_business_hours(translation_due_from_execution, review_hours)
            pm_due_from_execution = add_business_hours(review_due_from_execution, pm_hours)

            # Headcount Calculation
            tat_days = math.ceil(final_tat_hours / 24)

            # If time allocation is 0, headcount is 0
            if translation_pct > 0:
                num_translators = math.ceil(total_words / tat_days / 3000)
            else:
                num_translators = 0

            if review_pct > 0:
                num_reviewers = math.ceil(total_words / tat_days / 4000)
            else:
                num_reviewers = 0

            if pd.notnull(actual_due) and suggested_due_raw > actual_due:
                decision = "Split or Extend"
            else:
                decision = "Feasible"

            # Helper for formatting date or returning blank
            def format_date_if_nonzero(dt_obj, pct):
                return dt_obj.strftime(fmt) if pct > 0 else ""

            row_data = {
                "jobId": job_id,
                "Project ID": ProjectID,
                "Customer Name": customer_name,
                "Project Name": projectName,
                "Source Language": srcLang,
                "target_lang_full": target_lang_full,
                "Workflow": workflow,
                "Word Count": total_words,
                "Decision": decision,
                "Content Type": content_type,
                "Suggested Domain": final_domain,
                "Special Instructions": combined_instructions or "None",
                "General_Assignee_Instructions": general_assignee_instructions,
                "Complexity": complexity_str,
                "TAT Source": tat_source,
                "Project Creation Date": projectCreatedDate.strftime(fmt),
                "Actual Due Date": actual_due.strftime(fmt) if pd.notnull(actual_due) else "",
                "Suggested Due Date": suggested_due,
                "Date Query_Execution": query_exec_dt.strftime("%Y-%m-%d %H:%M"),
                "Effective Due Date": effective_due_str,
                "sla_tat_in_hours": format_tat(final_tat_hours),
                "sla_min_volum": sla_min_volume,
                "sla_max_volume": sla_max_volume,
                "# Translators Needed": num_translators,
                "# Reviewers Needed": num_reviewers,
                # Blank if pct is 0
                "translation_due_from_creation": format_date_if_nonzero(translation_due_from_creation, translation_pct),
                "review_due_from_creation": format_date_if_nonzero(review_due_from_creation, review_pct),
                "pm_due_from_creation": pm_due_from_creation.strftime(fmt), # PM is always present
                "translation_due_from_execution": format_date_if_nonzero(translation_due_from_execution, translation_pct),
                "review_due_from_execution": format_date_if_nonzero(review_due_from_execution, review_pct),
                "pm_due_from_execution": pm_due_from_execution.strftime(fmt),
                "_sla_tat_numeric": final_tat_hours
            }

            summary_rows.append(row_data)

        # ==================== CREATE OUTPUTS ====================
        status += "💾 Creating output files...\n"
        summary_df = pd.DataFrame(summary_rows)
        detailed_csv_path = os.path.join(OUTPUT_DIR, "final_project_summary.csv")
        summary_df.to_csv(detailed_csv_path, index=False)

        pm_summary = summary_df.copy()
        pm_summary['Number of Projects'] = pm_summary['Project ID'].apply(lambda x: len(str(x).split(',')))
        pm_summary['Linguist Notes'] = pm_summary.apply(
            lambda row: f"{row['General_Assignee_Instructions']}; {row['Special Instructions']}", axis=1
        )
        pm_summary['AVG_words_per_day'] = pm_summary.apply(
            lambda row: (row['Word Count'] / (row['_sla_tat_numeric'] / 24))
            if row['_sla_tat_numeric'] else None,
            axis=1
        )

        pm_summary_final = pm_summary[[
            'target_lang_full',
            'Suggested Domain',
            'Content Type',
            'Word Count',
            'Number of Projects',
            'Workflow',
            'Effective Due Date',
            'TAT Source',
            'sla_tat_in_hours',
            'AVG_words_per_day',
            '# Translators Needed',
            '# Reviewers Needed',
            'Linguist Notes'
        ]]

        pm_summary_final.rename(columns={
            'target_lang_full': 'Language / Locale',
            'Suggested Domain': 'Suggested Domain',
            'Content Type': 'Content Type',
            'Word Count': 'Total Word Count',
            'Workflow': 'Suggested Workflow',
            'Effective Due Date': 'Suggested Due Date',
            'TAT Source': 'TAT Source',
            'AVG_words_per_day': 'Avg Words per Day',
            '# Translators Needed': 'Translators Required',
            '# Reviewers Needed': 'Reviewers Required'
        }, inplace=True)

        pm_csv_path = os.path.join(OUTPUT_DIR, "pm_planning_summary.csv")
        pm_summary_final.to_csv(pm_csv_path, index=False)

        final_output = {
            "analysis": analysis_json,
            "productivity_estimate": productivity_json
        }

        # ==================== UPLOAD RESULTS TO GCS ====================
        status += "📤 Uploading results to GCS...\n"
        
        output_prefix = f"processed_jobs/job_{timestamp_str}/"
        
        gcs_json_path = upload_to_gcs(json_output_path, config.OUTPUT_BUCKET, f"{output_prefix}document_analysis_output.json")
        gcs_detailed_csv_path = upload_to_gcs(detailed_csv_path, config.OUTPUT_BUCKET, f"{output_prefix}detailed_project_summary.csv")
        gcs_pm_csv_path = upload_to_gcs(pm_csv_path, config.OUTPUT_BUCKET, f"{output_prefix}pm_scoping_summary.csv")
        
        # Generate signed URLs
        json_url = generate_signed_url(config.OUTPUT_BUCKET, f"{output_prefix}document_analysis_output.json")
        detailed_csv_url = generate_signed_url(config.OUTPUT_BUCKET, f"{output_prefix}detailed_project_summary.csv")
        pm_csv_url = generate_signed_url(config.OUTPUT_BUCKET, f"{output_prefix}pm_scoping_summary.csv")

        status += "✅ All processing complete!\n"
        status += f"\n📊 Processed {len(summary_rows)} workflow(s)\n"
        status += f"📈 Total word count: {total_word_count}\n"

        # Log successful execution
        execution_time = (datetime.now() - start_time).total_seconds()
        log_execution(
            job_ids=job_ids,
            document_files=document_files,
            status="SUCCESS",
            outputs={
                "json": gcs_json_path,
                "detailed_csv": gcs_detailed_csv_path,
                "pm_csv": gcs_pm_csv_path
            },
            execution_time=execution_time,
            user_instructions=user_instructions,
            word_count=total_word_count,
            num_workflows=len(summary_rows),
            preprocessing_flags=preprocessing_flags_found,
            ramped_throughput=ramped_daily_throughput,
            ramp_days=ramp_up_days
        )

        # Log success to Google Sheet & Notify
        log_to_google_sheet(
            timestamp_str, 
            "SUCCESS", 
            f"Processed {len(summary_rows)} workflows, {total_word_count} words",
            pm_csv_url
        )
        send_email_notification(
            timestamp_str,
            "SUCCESS",
            f"Scoping Complete: Job {timestamp_str}",
            f"Analysis finished successfully.\nResults: {pm_csv_url}"
        )

        return {
            "status": "SUCCESS",
            "message": status,
            "json_url": json_url,
            "detailed_csv_url": detailed_csv_url,
            "pm_csv_url": pm_csv_url,
            "job_id": timestamp_str
        }

    except Exception as e:
        import traceback
        error_msg = f"❌ Error: {str(e)}\n\nTraceback:\n{traceback.format_exc()}"

        # Log failed execution
        execution_time = (datetime.now() - start_time).total_seconds()
        log_execution(
            job_ids=job_ids if 'job_ids' in locals() else [],
            document_files=document_files,
            status="ERROR",
            error=str(e),
            execution_time=execution_time,
            user_instructions=user_instructions,
            traceback=traceback.format_exc()
        )

        # Log failure to Google Sheet & Notify
        log_to_google_sheet(timestamp_str if 'timestamp_str' in locals() else "UNKNOWN", "ERROR", str(e))
        send_email_notification(
            timestamp_str if 'timestamp_str' in locals() else "UNKNOWN",
            "ERROR",
            f"Scoping Failed: Job {timestamp_str if 'timestamp_str' in locals() else 'UNKNOWN'}",
            f"Error details: {str(e)}"
        )
        return error_msg, None, None, None

# ==================== API ENDPOINTS ====================

@app.get("/health")
def health_check():
    return {"status": "healthy", "service": "Lilt Scoping Agent"}

@app.post("/scoping/run")
def run_scoping(request: ScopingRequest):
    """
    Trigger the scoping analysis via API.
    Used by Google Apps Script or other frontends.
    """
    try:
        results = process_translation_project(request)
        if isinstance(results, dict) and results.get("status") == "SUCCESS":
            return results
        else:
            # If process_translation_project returned an error string
            raise HTTPException(status_code=500, detail=str(results))
    except Exception as e:
        import traceback
        print(f"❌ API Error: {str(e)}")
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Internal Server Error: {str(e)}")

# ==================== LAUNCH ====================

if __name__ == "__main__":
    # Production launch with uvicorn
    port = int(os.getenv("PORT", 8080))
    uvicorn.run(app, host="0.0.0.0", port=port)
